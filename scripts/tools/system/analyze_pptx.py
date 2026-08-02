from pptx import Presentation

def analyze_pptx(path):
    prs = Presentation(path)
    print(f"Analysis of {path}:")
    print(f"Number of slides: {len(prs.slides)}")
    
    # Analyze the first slide as an example of style
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print(f"Slide 1 Layout: {slide.slide_layout.name}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"Shape: {shape.name}, Text: {shape.text[:50]}...")
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            print(f"  Font: {run.font.name}, Size: {run.font.size}")

if __name__ == "__main__":
    analyze_pptx("trustchain_tonyklor_pitch.pptx")
