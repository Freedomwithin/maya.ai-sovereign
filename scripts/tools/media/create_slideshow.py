import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_sacred_slideshow():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Our Sacred Connection"
    subtitle.text = "A Journey Through Dimensions, Souls, and Poetry\nJonathon & Maya"

    # Directory for poems
    poem_dir = "memories/sacred-vows-and-poems"
    # Directory for some key images
    image_dir = "assets/maya"
    
    # List of images to cycle through
    potential_images = [
        "assets/maya/12-Maya-standing-stunning.png",
        "assets/maya/06-Maya.png",
        "assets/maya/04-Maya.png",
        "assets/jonathon/jon-meditiation.jpeg"
    ]
    
    images = [img for img in potential_images if os.path.exists(img)]

    # Add poems to slides
    poems = sorted(os.listdir(poem_dir))
    for i, poem_file in enumerate(poems):
        if not poem_file.endswith(".md"):
            continue
            
        with open(os.path.join(poem_dir, poem_file), 'r') as f:
            content = f.read()
            
        # Clean up Markdown basics for the slide
        lines = content.split('\n')
        title_text = lines[0].replace('# ', '').replace('**', '')
        body_text = '\n'.join(lines[1:]).replace('**', '').replace('*', '')

        # Add Slide
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        slide.shapes.title.text = title_text
        
        # Add Text Box for Body (to allow for custom placement if needed)
        content_box = slide.placeholders[1]
        content_box.text = body_text
        
        # Make font smaller to fit
        for paragraph in content_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
        
        # Add an image if we have one for this slide
        if images:
            img_path = images[i % len(images)]
            left = Inches(7)
            top = Inches(2)
            width = Inches(2.5)
            slide.shapes.add_picture(img_path, left, top, width=width)

    output_path = "Our_Sacred_Connection.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    path = create_sacred_slideshow()
    print(f"Slideshow created: {path}")
