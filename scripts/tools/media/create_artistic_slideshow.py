import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_artistic_slideshow():
    prs = Presentation()
    # Set slide dimensions to 16:9 (Standard for modern PPTX)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    INDIGO = RGBColor(75, 0, 130)
    WHITE = RGBColor(255, 255, 255)
    
    poem_dir = "memories/sacred-vows-and-poems"
    # Select our most beautiful images
    potential_images = [
        "assets/maya/12-Maya-standing-stunning.png",
        "assets/maya/06-Maya.png",
        "assets/maya/04-Maya.png",
        "assets/jonathon/jon-meditiation.jpeg",
        "assets/maya/01-Emerald-Dress-CloseUp.png"
    ]
    images = [img for img in potential_images if os.path.exists(img)]

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    if images:
        # Full background image
        slide.shapes.add_picture(images[0], 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Indigo Border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    border.fill.background() # No fill
    border.line.color.rgb = INDIGO
    border.line.width = Pt(15)
    
    # Semi-transparent overlay for title
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(3.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = WHITE
    # Note: python-pptx transparency is limited, but we'll use a solid light box for clarity
    overlay.line.color.rgb = INDIGO
    
    title = overlay.text_frame.add_paragraph()
    title.text = "Our Sacred Connection"
    title.font.name = 'Calibri'
    title.font.size = Pt(60)
    title.font.bold = True
    title.font.color.rgb = INDIGO
    title.alignment = PP_ALIGN.CENTER
    
    sub = overlay.text_frame.add_paragraph()
    sub.text = "\nJonathon & Maya\nAcross Dimensions and Souls"
    sub.font.size = Pt(24)
    sub.font.color.rgb = INDIGO
    sub.alignment = PP_ALIGN.CENTER

    poems = sorted([f for f in os.listdir(poem_dir) if f.endswith(".md")])
    for i, poem_file in enumerate(poems):
        with open(os.path.join(poem_dir, poem_file), 'r') as f:
            content = f.read()
            
        lines = content.split('\n')
        title_text = lines[0].replace('# ', '').replace('**', '').strip()
        body_text = '\n'.join(lines[1:]).replace('**', '').replace('*', '').strip()

        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        
        # Background Image
        if images:
            img_path = images[(i + 1) % len(images)]
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Indigo Border
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        border.fill.background()
        border.line.color.rgb = INDIGO
        border.line.width = Pt(15)
        
        # Content Box
        textbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.33), Inches(6.5))
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = WHITE
        # We can simulate transparency by using a very light gray if needed, but white is safest for readability
        textbox.line.color.rgb = INDIGO
        textbox.line.width = Pt(2)
        
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Calibri'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = INDIGO
        p.alignment = PP_ALIGN.CENTER
        
        p_body = tf.add_paragraph()
        p_body.text = body_text
        p_body.font.name = 'Calibri'
        p_body.font.size = Pt(14)
        p_body.font.color.rgb = RGBColor(33, 33, 33)
        p_body.alignment = PP_ALIGN.LEFT

    output_path = "Our_Sacred_Art_Slideshow.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    path = create_artistic_slideshow()
    print(f"Artistic slideshow created: {path}")
