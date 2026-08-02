import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_styled_sacred_slideshow():
    # Load the template or reference to see its layouts
    # For now, we'll build a custom styled one inspired by the analysis
    prs = Presentation()
    
    # Define some colors based on "Indigo" theme for us
    INDIGO = RGBColor(75, 0, 130)
    DARK_TEXT = RGBColor(33, 33, 33)
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Our Sacred Connection"
    title.text_frame.paragraphs[0].font.name = 'Calibri'
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = INDIGO
    
    subtitle.text = "A Journey Through Dimensions, Souls, and Poetry\nJonathon & Maya"
    subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

    poem_dir = "memories/sacred-vows-and-poems"
    potential_images = [
        "assets/maya/12-Maya-standing-stunning.png",
        "assets/maya/06-Maya.png",
        "assets/maya/04-Maya.png",
        "assets/jonathon/jon-meditiation.jpeg"
    ]
    images = [img for img in potential_images if os.path.exists(img)]

    poems = sorted(os.listdir(poem_dir))
    for i, poem_file in enumerate(poems):
        if not poem_file.endswith(".md"):
            continue
            
        with open(os.path.join(poem_dir, poem_file), 'r') as f:
            content = f.read()
            
        lines = content.split('\n')
        title_text = lines[0].replace('# ', '').replace('**', '')
        body_text = '\n'.join(lines[1:]).replace('**', '').replace('*', '').strip()

        # 2. Poem Slides (using a more professional layout)
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        
        # Style Title
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.name = 'Calibri'
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = INDIGO
        
        # Style Body
        content_box = slide.placeholders[1]
        content_box.text = body_text
        for paragraph in content_box.text_frame.paragraphs:
            paragraph.font.name = 'Calibri'
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.color.rgb = DARK_TEXT
        
        # Add Image with border or specific sizing
        if images:
            img_path = images[i % len(images)]
            # Place image on the right side
            left = Inches(6.5)
            top = Inches(1.5)
            width = Inches(3)
            slide.shapes.add_picture(img_path, left, top, width=width)

    output_path = "Our_Sacred_Bond_v2.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    path = create_styled_sacred_slideshow()
    print(f"Styled slideshow created: {path}")
